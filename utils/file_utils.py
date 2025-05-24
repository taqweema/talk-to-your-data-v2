import fitz  # PyMuPDF
from docx import Document
import pandas as pd
from pptx import Presentation

def read_txt(file):
    try:
        return file.read().decode("utf-8")
    except Exception as e:
        return f"[ERROR] Could not read TXT: {str(e)}"

def read_pdf(file):
    try:
        text = ""
        with fitz.open(stream=file.read(), filetype="pdf") as doc:
            for page in doc:
                text += page.get_text()
        return text
    except Exception as e:
        return f"[ERROR] Could not read PDF: {str(e)}"

def read_docx(file):
    try:
        doc = Document(file)
        return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
    except Exception as e:
        return f"[ERROR] Could not read DOCX: {str(e)}"

def read_csv(file):
    try:
        df = pd.read_csv(file)
        return df.to_string(index=False)
    except Exception as e:
        return f"[ERROR] Could not read CSV: {str(e)}"

def read_excel(file):
    try:
        df = pd.read_excel(file)
        return df.to_string(index=False)
    except Exception as e:
        return f"[ERROR] Could not read Excel: {str(e)}"

def read_pptx(file):
    try:
        prs = Presentation(file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        return f"[ERROR] Could not read PPTX: {str(e)}"

def read_file(file):
    filename = file.name.lower()
    if filename.endswith(".txt"):
        return read_txt(file)
    elif filename.endswith(".pdf"):
        return read_pdf(file)
    elif filename.endswith(".docx"):
        return read_docx(file)
    elif filename.endswith(".csv"):
        return read_csv(file)
    elif filename.endswith(".xls") or filename.endswith(".xlsx"):
        return read_excel(file)
    elif filename.endswith(".pptx"):
        return read_pptx(file)
    else:
        return "[ERROR] Unsupported file type"
